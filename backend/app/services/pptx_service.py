from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from pdf2image import convert_from_path

class PPTXService:
    def __init__(self):
        pass

    def create_pptx_from_ocr(self, pdf_path: str, ocr_result: dict, output_path: str):
        """
        Creates a PPTX file using OCR data and PDF images as background.
        Optimized for memory: Processes one page at a time.
        """
        from pdf2image import pdfinfo_from_path, convert_from_path
        
        prs = Presentation()
        
        # 1. Get total pages without converting (Lightweight)
        try:
            info = pdfinfo_from_path(pdf_path)
            total_pages = info["Pages"]
        except Exception as e:
            print(f"Error reading PDF info: {e}")
            raise e

        # Group elements by page
        elements = ocr_result.get("elements", [])
        pages_elements = {}
        for el in elements:
            page_idx = el.get("page", 1) - 1 # API is 1-indexed usually
            if page_idx not in pages_elements:
                pages_elements[page_idx] = []
            pages_elements[page_idx].append(el)

        # 2. Process each page iteratively
        for i in range(total_pages):
            # Convert SINGLE page to image (Memory Efficient)
            # 150 DPI is balanced for screen viewing
            try:
                images = convert_from_path(pdf_path, dpi=150, first_page=i+1, last_page=i+1)
                if not images:
                    continue
                image = images[0]
            except Exception as e:
                print(f"Error converting page {i+1}: {e}")
                continue

            # Set slide size based on first page (only once)
            if i == 0:
                width, height = image.size
                dpi = 150
                prs.slide_width = int(width / dpi * 914400)
                prs.slide_height = int(height / dpi * 914400)

            # Save temp image as JPEG (Smaller size than PNG)
            temp_img_path = f"temp_page_{i}.jpg"
            image.save(temp_img_path, "JPEG", quality=80)

            # Create slide
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            # Add background image
            slide.shapes.add_picture(temp_img_path, 0, 0, prs.slide_width, prs.slide_height)
            
            # Add text boxes and accumulate text for notes
            page_elements = pages_elements.get(i, [])
            all_page_text = []
            for el in page_elements:
                # Need original image dimensions for scaling logic inside this method
                extracted_text = self._add_element_to_slide(slide, el, prs.slide_width, prs.slide_height, width, height)
                if extracted_text:
                    all_page_text.append(extracted_text)

            # Add all extracted text to slide notes (fallback)
            if all_page_text:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = "\n\n".join(all_page_text)

            # Cleanup immediately
            if os.path.exists(temp_img_path):
                os.remove(temp_img_path)
            
            # Explicitly delete image object to free RAM
            del image
            del images

        prs.save(output_path)
        return output_path

    def _add_element_to_slide(self, slide, element, slide_w_emu, slide_h_emu, img_w_px, img_h_px):
        import re
        category = element.get("category", "")
        # Filter unwanted categories (but NOT figure - we extract alt text from figures)
        if category in ["chart", "header", "footer"]:
            return None

        content = element.get("content", {})
        text = content.get("text", "")
        html = content.get("html", "")
        
        # For figures, extract text from alt attribute
        if category == "figure":
            # Try to find alt attribute with meaningful text
            alt_match = re.search(r'alt="([^"]+)"', html)
            if alt_match:
                alt_text = alt_match.group(1)
                # Skip if alt is empty or just special characters
                if alt_text and len(alt_text.strip()) > 2 and not all(c in '□●◇◆↓↑←→' for c in alt_text.strip()):
                    text = alt_text.replace('\\n', '\n')
                else:
                    return None
            else:
                return None
        
        # If text is empty, try to extract from HTML (for non-figure elements)
        if not text and html:
            # Replace <br>, </tr>, </td> with newline to preserve structure
            text = re.sub(r'</?(tr|td|div|p)\s*[^>]*>', '\n', html)
            text = re.sub(r'<br\s*/?>', '\n', html, flags=re.IGNORECASE)
            # Remove remaining HTML tags
            text = re.sub(r'<[^>]+>', '', text).strip()
            # Clean up multiple newlines
            text = re.sub(r'\n\s*\n', '\n', text).strip()
            
        if not text:
            return None

        # Coordinates are normalized (0-1) in 'coordinates' list
        coords = element.get("coordinates", [])
        if not coords:
            return

        # Get min/max from normalized coordinates
        xs = [p['x'] for p in coords]
        ys = [p['y'] for p in coords]
        
        if not xs or not ys:
            return
            
        x_min, x_max = min(xs), max(xs)
        y_min, y_max = min(ys), max(ys)
        
        # Convert to EMU
        x_emu = int(x_min * slide_w_emu)
        y_emu = int(y_min * slide_h_emu)
        w_emu = int((x_max - x_min) * slide_w_emu)
        h_emu = int((y_max - y_min) * slide_h_emu)
        
        # Add padding to width to prevent aggressive wrapping (approx 15% buffer)
        w_emu = int(w_emu * 1.15)
        
        # Ensure minimum size to be visible
        if w_emu < 100 or h_emu < 100:
            return

        # Create textbox
        textbox = slide.shapes.add_textbox(x_emu, y_emu, w_emu, h_emu)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.text = text
        
        # Styling
        # Try to extract font size from HTML style style='font-size:22px'
        font_size_pt = 12 # default
        font_match = re.search(r"font-size:(\d+)px", html)
        if font_match:
            # HTML pixels to PPT points. Scaling down slightly helps fit text.
            px_size = int(font_match.group(1))
            font_size_pt = px_size * 0.65
            
        if tf.paragraphs:
            p = tf.paragraphs[0]
            # Ensure minimum font size
            p.font.size = Pt(max(9, font_size_pt))
            # Optional: color, bold based on tag (h1 vs p)
            if category.startswith("heading"):
                p.font.bold = True
        
        # Make background white and transparent-ish? 
        # User requested "editable text boxes". 
        # To mask the underlying baked-in text in the image, we likely need a solid background.
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        return text
